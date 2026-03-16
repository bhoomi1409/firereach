"""
PowerPoint Generation Service
Creates personalized pitch decks for each target company
"""

import os
import io
from datetime import datetime
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import httpx
import asyncio

class PPTService:
    def __init__(self):
        self.template_path = "templates/firereach_template.pptx"
        self.output_dir = "generated_ppts"
        os.makedirs(self.output_dir, exist_ok=True)
    
    async def generate_personalized_ppt(
        self,
        company_data: Dict,
        contact_data: Dict,
        signals: List[Dict],
        icp_data: Dict
    ) -> str:
        """
        Generate a personalized PowerPoint presentation for the target company
        Returns the file path of the generated PPT
        """
        company_name = company_data.get("name", "Company")
        
        # Create new presentation
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Slide 1: Title Slide
        self._create_title_slide(prs, company_name, contact_data, icp_data)
        
        # Slide 2: Problem Statement (Personalized)
        self._create_problem_slide(prs, company_data, signals, icp_data)
        
        # Slide 3: Solution Overview
        self._create_solution_slide(prs, icp_data)
        
        # Slide 4: Why Now? (Company-specific signals)
        self._create_timing_slide(prs, company_name, signals)
        
        # Slide 5: Social Proof & Results
        self._create_social_proof_slide(prs, company_data)
        
        # Slide 6: ROI Calculator (Personalized)
        self._create_roi_slide(prs, company_data, icp_data)
        
        # Slide 7: Next Steps
        self._create_next_steps_slide(prs, contact_data)
        
        # Save presentation
        filename = f"{company_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        
        return filepath
    
    def _create_title_slide(self, prs: Presentation, company_name: str, contact_data: Dict, icp_data: Dict):
        """Create personalized title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)  # Dark slate
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = f"Scaling {company_name}'s Outreach"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(251, 146, 60)  # Orange
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"AI-Powered Outreach Automation for {contact_data.get('first_name', 'Your Team')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(148, 163, 184)  # Gray
        
        # Date and personalization
        date_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(14), Inches(1))
        date_frame = date_box.text_frame
        date_frame.text = f"Prepared for {contact_data.get('first_name', '')} {contact_data.get('last_name', '')} • {datetime.now().strftime('%B %Y')}"
        date_para = date_frame.paragraphs[0]
        date_para.alignment = PP_ALIGN.CENTER
        date_para.font.size = Pt(16)
        date_para.font.color.rgb = RGBColor(100, 116, 139)
    
    def _create_problem_slide(self, prs: Presentation, company_data: Dict, signals: List[Dict], icp_data: Dict):
        """Create problem statement slide with company-specific context"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"The Challenge: {company_data.get('name', 'Your Company')}'s Growth Demands"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(251, 146, 60)
        
        # Problem points based on signals and ICP
        problems = self._generate_problem_points(company_data, signals, icp_data)
        
        y_position = 2
        for i, problem in enumerate(problems):
            problem_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_position), Inches(13), Inches(1.2))
            problem_frame = problem_box.text_frame
            problem_frame.text = f"• {problem}"
            problem_para = problem_frame.paragraphs[0]
            problem_para.font.size = Pt(20)
            problem_para.font.color.rgb = RGBColor(226, 232, 240)
            y_position += 1.3
    
    def _create_solution_slide(self, prs: Presentation, icp_data: Dict):
        """Create solution overview slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "FireReach: Autonomous Outreach Engine"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(251, 146, 60)
        
        # Solution features
        features = [
            "🎯 AI discovers ideal prospects from your ICP",
            "📊 Enriches companies with live signals (funding, hiring, news)",
            "🔍 Finds decision-maker contacts automatically",
            "✍️ Generates hyper-personalized emails using AI",
            "📧 Sends at scale with full compliance",
            "📈 Tracks performance and optimizes continuously"
        ]
        
        y_position = 2.5
        for feature in features:
            feature_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_position), Inches(13), Inches(0.8))
            feature_frame = feature_box.text_frame
            feature_frame.text = feature
            feature_para = feature_frame.paragraphs[0]
            feature_para.font.size = Pt(18)
            feature_para.font.color.rgb = RGBColor(226, 232, 240)
            y_position += 0.9
    
    def _create_timing_slide(self, prs: Presentation, company_name: str, signals: List[Dict]):
        """Create 'Why Now?' slide with company-specific signals"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"Why Now? {company_name}'s Growth Signals"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(251, 146, 60)
        
        # Display top signals
        y_position = 2.5
        for i, signal in enumerate(signals[:4]):  # Top 4 signals
            signal_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_position), Inches(13), Inches(1))
            signal_frame = signal_box.text_frame
            
            # Add signal type emoji
            emoji = "💰" if "funding" in signal.get("type", "") else \
                   "👥" if "hiring" in signal.get("type", "") else \
                   "📰" if "news" in signal.get("type", "") else "🚀"
            
            signal_frame.text = f"{emoji} {signal.get('summary', '')[:100]}..."
            signal_para = signal_frame.paragraphs[0]
            signal_para.font.size = Pt(16)
            signal_para.font.color.rgb = RGBColor(226, 232, 240)
            y_position += 1.2
    
    def _create_social_proof_slide(self, prs: Presentation, company_data: Dict):
        """Create social proof slide with relevant case studies"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Results from Similar Companies"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(251, 146, 60)
        
        # Case studies (would be dynamically selected based on company industry/stage)
        case_studies = [
            {
                "company": "TechCorp (Series B SaaS)",
                "result": "3x reply rates, 50% more qualified meetings",
                "metric": "From 2% to 6% reply rate in 30 days"
            },
            {
                "company": "DataFlow (B2B Platform)", 
                "result": "Reduced SDR workload by 80%",
                "metric": "$2M pipeline generated in Q1"
            },
            {
                "company": "CloudScale (Enterprise SaaS)",
                "result": "Automated 90% of outbound process",
                "metric": "200% increase in outbound velocity"
            }
        ]
        
        y_position = 2.5
        for case in case_studies:
            # Company name
            company_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_position), Inches(13), Inches(0.6))
            company_frame = company_box.text_frame
            company_frame.text = case["company"]
            company_para = company_frame.paragraphs[0]
            company_para.font.size = Pt(18)
            company_para.font.bold = True
            company_para.font.color.rgb = RGBColor(251, 146, 60)
            
            # Result
            result_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_position + 0.4), Inches(13), Inches(0.8))
            result_frame = result_box.text_frame
            result_frame.text = f"✓ {case['result']}\n  {case['metric']}"
            result_para = result_frame.paragraphs[0]
            result_para.font.size = Pt(14)
            result_para.font.color.rgb = RGBColor(226, 232, 240)
            
            y_position += 1.5
    
    def _create_roi_slide(self, prs: Presentation, company_data: Dict, icp_data: Dict):
        """Create ROI calculator slide with company-specific numbers"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"ROI Projection for {company_data.get('name', 'Your Company')}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(251, 146, 60)
        
        # ROI calculation (simplified)
        estimated_headcount = company_data.get("headcount", 100)
        sdr_count = max(2, estimated_headcount // 50)  # Estimate SDR team size
        
        roi_data = [
            f"Current SDR team: ~{sdr_count} people",
            f"Manual outreach capacity: ~{sdr_count * 50} emails/month",
            f"FireReach capacity: ~{sdr_count * 500} emails/month (10x)",
            f"Expected reply rate improvement: 2% → 6% (3x)",
            f"Additional qualified meetings: +{sdr_count * 20}/month",
            f"Estimated pipeline value: ${sdr_count * 100}K/month"
        ]
        
        y_position = 2.5
        for data_point in roi_data:
            roi_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_position), Inches(13), Inches(0.7))
            roi_frame = roi_box.text_frame
            roi_frame.text = f"• {data_point}"
            roi_para = roi_frame.paragraphs[0]
            roi_para.font.size = Pt(16)
            roi_para.font.color.rgb = RGBColor(226, 232, 240)
            y_position += 0.8
    
    def _create_next_steps_slide(self, prs: Presentation, contact_data: Dict):
        """Create next steps slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Ready to Scale Your Outreach?"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(42)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(251, 146, 60)
        
        # Next steps
        steps_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(4))
        steps_frame = steps_box.text_frame
        steps_text = f"""Let's discuss how FireReach can help {contact_data.get('first_name', 'your team')} achieve similar results.

📅 Book a 15-minute demo
🎯 See FireReach in action with your ICP
📊 Get a custom ROI analysis
🚀 Start your pilot program

Ready when you are, {contact_data.get('first_name', '')}!"""
        
        steps_frame.text = steps_text
        steps_para = steps_frame.paragraphs[0]
        steps_para.alignment = PP_ALIGN.CENTER
        steps_para.font.size = Pt(18)
        steps_para.font.color.rgb = RGBColor(226, 232, 240)
    
    def _generate_problem_points(self, company_data: Dict, signals: List[Dict], icp_data: Dict) -> List[str]:
        """Generate company-specific problem points based on signals and data"""
        problems = []
        
        # Base problems from ICP
        base_problems = [
            "Manual outreach processes limit scale and personalization",
            "Low reply rates from generic, template-based emails", 
            "SDR time wasted on research instead of selling",
            "Difficulty identifying the right prospects at the right time"
        ]
        
        # Add signal-based problems
        for signal in signals:
            signal_type = signal.get("type", "")
            if "funding" in signal_type:
                problems.append("Recent funding means pressure to scale revenue quickly")
            elif "hiring" in signal_type:
                problems.append("Rapid hiring creates need for efficient lead generation")
            elif "expansion" in signal.get("summary", "").lower():
                problems.append("Market expansion requires targeted outreach at scale")
        
        # Combine and limit to top 4
        all_problems = base_problems + problems
        return all_problems[:4]

ppt_service = PPTService()